"""Load text from supported document formats (with optional PaddleOCR)."""

from __future__ import annotations

from pathlib import Path

from infrarag.config import OcrConfig
from infrarag.ingest.ocr import (
    IMAGE_EXTENSIONS,
    OcrUnavailableError,
    ocr_image,
    ocr_pil_images,
    pdf_text_needs_ocr,
    render_pdf_pages,
)
from infrarag.models import Document

_DEFAULT_OCR = OcrConfig(
    enabled=True,
    lang="es",
    use_gpu=False,
    min_text_chars=40,
    max_pdf_pages=50,
    skip_photo_folders=True,
    photo_folder_image_ratio=0.6,
    photo_folder_min_files=5,
)


def load_document(path: Path, *, ocr: OcrConfig | None = None) -> Document:
    """Extract plain text from a supported file path."""
    ocr_config = ocr or _DEFAULT_OCR
    path = path.resolve()
    suffix = path.suffix.lower()
    used_ocr = False
    truncated = False

    if suffix in {".txt", ".md"}:
        text = path.read_text(encoding="utf-8", errors="replace")
    elif suffix == ".pdf":
        text, used_ocr, truncated = _load_pdf(path, ocr_config)
    elif suffix in IMAGE_EXTENSIONS:
        text = _load_image(path, ocr_config)
        used_ocr = True
    elif suffix == ".docx":
        text = _load_docx(path)
    elif suffix == ".xlsx":
        text = _load_xlsx(path)
    elif suffix == ".pptx":
        text = _load_pptx(path)
    elif suffix in {".odt", ".ods", ".odp"}:
        text = _load_odf(path)
    else:
        raise ValueError(f"Unsupported extension: {suffix}")

    text = text.strip()
    if not text:
        raise ValueError(f"No extractable text in {path}")

    metadata = {"suffix": suffix}
    if used_ocr:
        metadata["ocr"] = "true"
    if truncated:
        metadata["ocr_truncated"] = "true"
    return Document(source_path=str(path), text=text, metadata=metadata)


def _load_image(path: Path, ocr: OcrConfig) -> str:
    if not ocr.enabled:
        raise OcrUnavailableError(
            f"Image {path} requires OCR but ingest.ocr.enabled is false"
        )
    return ocr_image(path, ocr)


def _load_pdf(path: Path, ocr: OcrConfig) -> tuple[str, bool, bool]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts: list[str] = []
    for page in reader.pages:
        parts.append(page.extract_text() or "")
    text = "\n".join(parts)
    page_count = len(reader.pages)

    if not ocr.enabled or not pdf_text_needs_ocr(
        text, page_count=page_count, min_text_chars=ocr.min_text_chars
    ):
        return text, False, False

    images = render_pdf_pages(path, max_pages=ocr.max_pdf_pages)
    truncated = page_count > ocr.max_pdf_pages
    ocr_text = ocr_pil_images(images, ocr)
    if not ocr_text:
        raise ValueError(f"OCR produced no text for scanned PDF {path}")
    return ocr_text, True, truncated


def _load_docx(path: Path) -> str:
    from docx import Document as DocxDocument

    doc = DocxDocument(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def _load_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    lines: list[str] = []
    for sheet in wb.worksheets:
        lines.append(f"# Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if c is None else str(c) for c in row]
            if any(cells):
                lines.append("\t".join(cells))
    return "\n".join(lines)


def _load_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)


def _load_odf(path: Path) -> str:
    from odf import teletype, text
    from odf.opendocument import load

    doc = load(str(path))
    paragraphs = doc.getElementsByType(text.P)
    return "\n".join(teletype.extractText(p) for p in paragraphs)
